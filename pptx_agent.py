"""
このモジュールはパワーポイント資料を自動生成するAIエージェントを実装します。
LangGraphを使ってワークフロー（グラフ）を定義し、複数のLLMを連携させて処理を実行します。
APIクォータ超過時のフォールバックメカニズムを備えています。

Classes:
    PPTXAgent: PowerPoint資料を自動生成するAIエージェント
"""

from typing import Any, Optional, Dict, Callable
import os
import logging
import time
from langchain_openai import ChatOpenAI
from langgraph.graph import END, StateGraph

from datamodel import Judgement, State
from story_generator import StoryGenerator
from story_evaluator import StoryEvaluator
from slide_contents_generator import SlideContentsGenerator
from pptx_code_generator import PPTXCodeGenerator

# ロガーの設定
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PPTXAgent:
    """
    PowerPoint資料を自動生成するAIエージェント

    Attributes:
        story_generator (StoryGenerator): ストーリー生成器
        story_evaluator (StoryEvaluator): ストーリー評価器
        slide_contents_generator (SlideContentsGenerator): スライド内容生成器
        pptx_code_generator (PPTXCodeGenerator): PPTXコード生成器
        graph (StateGraph): ワークフローを表すグラフ
        use_fallback (bool): APIエラー時にフォールバックを使用するかどうか
        fallback_llm (ChatOpenAI): フォールバック用のLLM
        max_retries (int): 最大再試行回数
    """
    def __init__(self, llm: ChatOpenAI, use_fallback: bool = True, max_retries: int = 3):
        """
        PPTXAgentクラスの初期化

        Args:
            llm (ChatOpenAI): 対話型言語モデルのインスタンス
            use_fallback (bool, optional): APIエラー時にフォールバックを使用するかどうか。デフォルトはTrue
            max_retries (int, optional): 最大再試行回数。デフォルトは3
        """
        # フォールバックの設定
        self.use_fallback = use_fallback
        self.primary_llm = llm
        self.fallback_llm = None
        if self.use_fallback and llm.model_name != "gpt-3.5-turbo":
            try:
                self.fallback_llm = ChatOpenAI(model="gpt-3.5-turbo", temperature=0.0)
            except Exception as e:
                logger.warning(f"フォールバックLLMの初期化に失敗しました: {e}")
                
        self.max_retries = max_retries
        
        # 各種ジェネレーターの初期化
        self.story_generator = StoryGenerator(llm=llm)
        self.story_evaluator = StoryEvaluator(llm=llm)
        self.slide_contents_generator = SlideContentsGenerator(llm=llm)
        self.pptx_code_generator = PPTXCodeGenerator(llm=llm)
        
        # グラフの作成
        self.graph = self._create_graph()
        
    def _with_retries_and_fallback(self, func: Callable, *args, **kwargs) -> Any:
        """
        リトライとフォールバックを実装した関数ラッパー

        Args:
            func (Callable): 実行する関数
            *args: 関数に渡す位置引数
            **kwargs: 関数に渡すキーワード引数

        Returns:
            Any: 関数の実行結果

        Raises:
            Exception: すべての試行が失敗した場合
        """
        last_error = None
        
        # 通常のLLMで試行
        for attempt in range(self.max_retries):
            try:
                return func(*args, **kwargs)
            except Exception as e:
                last_error = e
                error_msg = str(e).lower()
                
                # APIクォータ超過エラーの場合はすぐにフォールバック
                if "insufficient_quota" in error_msg or "rate_limit" in error_msg:
                    logger.warning(f"APIクォータ超過またはレート制限エラー: {e}")
                    break
                
                # その他のエラーの場合はリトライ
                logger.warning(f"エラーが発生しました（試行 {attempt+1}/{self.max_retries}）: {e}")
                time.sleep(2 ** attempt)  # 指数バックオフ
        
        # フォールバックが有効で利用可能な場合
        if self.use_fallback and self.fallback_llm is not None:
            logger.info("フォールバックLLMを使用します")
            
            # 元のLLMを一時的にフォールバックに置き換え
            orig_llm = {}
            for component_name in ["story_generator", "story_evaluator", "slide_contents_generator", "pptx_code_generator"]:
                component = getattr(self, component_name)
                orig_llm[component_name] = component.llm
                component.llm = self.fallback_llm
            
            try:
                # フォールバックLLMで再実行
                return func(*args, **kwargs)
            except Exception as e:
                logger.error(f"フォールバックLLMでも失敗しました: {e}")
                last_error = e
            finally:
                # 元のLLMを復元
                for component_name, llm in orig_llm.items():
                    component = getattr(self, component_name)
                    component.llm = llm
        
        # すべての試行が失敗した場合
        if last_error:
            raise last_error
        
    def _create_graph(self) -> StateGraph:
        """
        ワークフローグラフを作成する

        Returns:
            StateGraph: コンパイル済みのワークフローグラフ
        """
        # グラフの初期化
        workflow = StateGraph(State)
        
        # 各ノードの追加
        workflow.add_node("generate_story", self._generate_story)
        workflow.add_node("evaluate_story", self._evaluate_story)
        workflow.add_node("generate_slide_contents", self._generate_slide_contents)
        workflow.add_node("generate_pptx_code", self._generate_pptx_code)
        
        # エントリーポイントの設定
        workflow.set_entry_point("generate_story")
        
        # ノード間のエッジの追加
        workflow.add_edge("generate_story", "evaluate_story")
        workflow.add_conditional_edges(
            "evaluate_story",
            lambda state: not state.current_judge and state.iteration < 5,
            {True: "generate_story", False: "generate_slide_contents"}
        )
        workflow.add_edge("generate_slide_contents", "generate_pptx_code")
        workflow.add_edge("generate_pptx_code", END)
        
        # グラフのコンパイル
        return workflow.compile()
    
    def _generate_story(self, state: State) -> dict[str, Any]:
        """
        ストーリーを生成するノード処理

        Args:
            state (State): 現在の状態

        Returns:
            dict[str, Any]: 更新する状態の要素
        """
        # ストーリーの生成（リトライとフォールバックを使用）
        def generate():
            return self.story_generator.run(state.user_request)
        
        new_story = self._with_retries_and_fallback(generate)
        
        return {
            "story": new_story,
            "iteration": state.iteration + 1
        }
        
    def _evaluate_story(self, state: State) -> dict[str, Any]:
        """
        ストーリーを評価するノード処理

        Args:
            state (State): 現在の状態

        Returns:
            dict[str, Any]: 更新する状態の要素
        """
        # ストーリーの評価（リトライとフォールバックを使用）
        def evaluate():
            return self.story_evaluator.run(state.user_request, state.story)
        
        judgement = self._with_retries_and_fallback(evaluate)
        
        return {
            "current_judge": judgement.judge,
            "judgement_reason": judgement.reason
        }
        
    def _generate_slide_contents(self, state: State) -> dict[str, Any]:
        """
        スライド内容を生成するノード処理

        Args:
            state (State): 現在の状態

        Returns:
            dict[str, Any]: 更新する状態の要素
        """
        # スライド内容の生成（リトライとフォールバックを使用）
        def generate_contents():
            return self.slide_contents_generator.run(state.user_request, state.story)
        
        slide_contents = self._with_retries_and_fallback(generate_contents)
        
        return {"slide_contents": slide_contents}
    
    def _generate_pptx_code(self, state: State) -> dict[str, Any]:
        """
        PPTXコードを生成するノード処理

        Args:
            state (State): 現在の状態

        Returns:
            dict[str, Any]: 更新する状態の要素
        """
        # Python-pptxコードの生成（リトライとフォールバックを使用）
        def generate_code():
            return self.pptx_code_generator.run(state.slide_contents)
        
        pptx_code = self._with_retries_and_fallback(generate_code)
        
        return {"slide_gen_code": pptx_code}
    
    def run(self, user_request: str) -> str:
        """
        AIエージェントを実行する

        Args:
            user_request (str): ユーザーからのリクエスト

        Returns:
            str: 生成されたPythonコード
        """
        # 初期状態の設定
        initial_state = State(user_request=user_request)
        
        # 出力ディレクトリの確認
        os.makedirs("workspace/output", exist_ok=True)
        
        try:
            # グラフの実行
            final_state = self.graph.invoke(initial_state)
            # 最終的なPythonコードの取得
            return final_state["slide_gen_code"]
        except Exception as e:
            logger.error(f"グラフの実行中にエラーが発生しました: {e}")
            # エラーが発生した場合は、簡易的なコードテンプレートを返す
            return """
# エラーが発生したため、簡易的なPPTXファイルを生成します
from pptx import Presentation
from pptx.util import Inches, Pt

# 新しいプレゼンテーションを作成
prs = Presentation()

# タイトルスライドを追加
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "エラーが発生しました"
subtitle.text = "APIクォータ超過またはその他のエラーにより、プレゼンテーションを生成できませんでした。"

# 情報スライドを追加
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]

title.text = "エラーの解決方法"
tf = body.text_frame
tf.text = "以下の方法を試してください："

p = tf.add_paragraph()
p.text = "1. OpenAIダッシュボードでAPIの使用状況と制限を確認する"
p.level = 1

p = tf.add_paragraph()
p.text = "2. 有料プランにアップグレードする"
p.level = 1

p = tf.add_paragraph()
p.text = "3. 別のAPIキーを使用する"
p.level = 1

# プレゼンテーションを保存
prs.save('workspace/output/error_presentation.pptx')
print("エラー用のプレゼンテーションが生成されました: workspace/output/error_presentation.pptx")
""" 