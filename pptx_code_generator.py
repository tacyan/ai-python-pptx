"""
このモジュールはスライド内容をPythonコードに変換するための機能を提供します。
生成されたスライド内容を、python-pptxを使用してPowerPointファイルを生成するためのコードに変換します。
画像や動画も挿入できるように拡張されています。
OpenAIとGoogle Gemini両方のAPIに対応しています。

Classes:
    PPTXCodeGenerator: PowerPointスライド生成コードを生成するクラス
"""

import os
import logging
from typing import Optional
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.language_models.chat_models import BaseChatModel

# ロガーの設定
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Gemini APIの利用可能性をチェック
try:
    from langchain_google_genai import ChatGoogleGenerativeAI
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False
    logger.warning("langchain_google_genai をインポートできませんでした。Google Gemini機能は無効化されます。")

class PPTXCodeGenerator:
    """
    PowerPointスライド生成コードを生成するクラス

    Attributes:
        llm (BaseChatModel): 対話型言語モデルのインスタンス
        max_retries (int): APIコール失敗時の最大再試行回数
    """
    def __init__(self, llm: BaseChatModel, max_retries: int = 2):
        """
        PPTXCodeGeneratorクラスの初期化

        Args:
            llm (BaseChatModel): 対話型言語モデルのインスタンス
            max_retries (int, optional): 最大再試行回数。デフォルトは2
        """
        self.llm = llm
        self.max_retries = max_retries
        # APIプロバイダーの検出（クラス名でチェック）
        self.api_provider = self._detect_api_provider(llm)
        logger.info(f"PPTXCodeGenerator initialized with {self.api_provider} API")
        
        # 画像アップロードディレクトリの確保
        os.makedirs("workspace/input/images", exist_ok=True)
        
    def _detect_api_provider(self, llm: BaseChatModel) -> str:
        """
        使用されているAPIプロバイダーを検出する
        
        Args:
            llm (BaseChatModel): 対話型言語モデルのインスタンス
            
        Returns:
            str: 検出されたAPIプロバイダー名 ("OpenAI" または "Google Gemini")
        """
        llm_class_str = str(llm.__class__).lower()
        
        if "openai" in llm_class_str:
            return "OpenAI"
        elif "google" in llm_class_str or "gemini" in llm_class_str:
            return "Google Gemini"
        else:
            logger.warning(f"不明なLLMタイプです: {llm.__class__.__name__}")
            return "Unknown"
        
    def run(self, slide_contents: str) -> str:
        """
        スライド内容からPowerPointスライド生成コードを生成する

        Args:
            slide_contents (str): 生成されたスライド内容

        Returns:
            str: PowerPointスライド生成のためのPythonコード
        """
        # プロンプトを定義
        prompt = ChatPromptTemplate.from_messages(
            [
                (
                    "system",
                    "python-pptxモジュールを用いてプレゼンテーション資料のスライドを自動生成する専門家です。"
                ),
                (
                    "human",
                    "以下のスライドの内容を生成するためのpythonコードを生成してください。\n\n"
                    "スライドの内容:\n{slide_contents}\n\n"
                    "以下のpptxファイルを読み込み、テンプレートとして使用してください。\n"
                    "workspace/input/template.pptx\n"
                    "テンプレートのレイアウト情報は以下を参照してください。記載にないレイアウト番号やプレースホルダー番号は決して使用しないでください。\n"
                    "    タイトルスライド: slide_layouts[2]\n"
                    "        placeholder_format.idx:\n"
                    "            0: 会社名など\n"
                    "            10: 発表タイトル\n"
                    "            11: サブタイトル・日付など\n"
                    "            12: 発表者名など\n"
                    "    一般スライド: slide_layouts[0]\n"
                    "        placeholder_format.idx:\n"
                    "            0: スライドタイトル\n"
                    "            1: 内容\n"
                    "作成したパワーポイントはworkspace/output内に出力されるようにしてください。\n\n"
                    "ルール:\n"
                    "- 【重要】必ずpython-pptxモジュールを使用したpythonコードのみを出力してください。\n"
                    "- 使用が許可されているのは、テキスト、図形、表、画像、動画です。\n"
                    "- テキスト以外の要素（図形、表、画像、動画）を使用してほしい箇所には、その旨が明記されています。\n"
                    "- 画像を挿入する場合は `workspace/input/images/` ディレクトリの画像ファイルを使用するコードを生成してください。\n"
                    "- 動画を挿入する場合も同様に `workspace/input/images/` ディレクトリのファイルを使用するコードを生成してください。\n"
                    "- '---next---' はスライド番号を進める合図です。このタイミングで新たなスライドを追加してください。\n\n"
                    "### 画像や動画の挿入方法について ###\n"
                    "以下のようなコード例を参考にしてください。\n"
                    "画像の挿入例:\n"
                    "```python\n"
                    "from pptx.util import Inches\n"
                    "slide = prs.slides.add_slide(prs.slide_layouts[0])\n"
                    "title = slide.shapes.title\n"
                    "title.text = \"画像の例\"\n"
                    "# 画像を追加\n"
                    "img_path = \"workspace/input/images/example.jpg\"\n"
                    "left = Inches(1)\n"
                    "top = Inches(2.5)\n"
                    "width = Inches(5)\n"
                    "slide.shapes.add_picture(img_path, left, top, width=width)\n"
                    "```\n"
                    "動画の挿入例:\n"
                    "```python\n"
                    "from pptx.util import Inches\n"
                    "slide = prs.slides.add_slide(prs.slide_layouts[0])\n"
                    "title = slide.shapes.title\n"
                    "title.text = \"動画の例\"\n"
                    "# 動画を追加\n"
                    "movie_path = \"workspace/input/images/example.mp4\"\n"
                    "left = Inches(2)\n"
                    "top = Inches(2.5)\n"
                    "width = Inches(5)\n"
                    "height = Inches(3)\n"
                    "slide.shapes.add_movie(movie_path, left, top, width, height)\n"
                    "```\n"
                )
            ]
        )
        
        # エラーハンドリングを追加
        last_error = None
        for attempt in range(self.max_retries + 1):
            try:
                # スライド生成のためのチェーンを作成
                chain = prompt | self.llm | StrOutputParser()
                # スライド生成のコードを生成
                result = chain.invoke({"slide_contents": slide_contents})
                logger.info(f"{self.api_provider} APIでPPTXコード生成に成功しました")
                return result
            except Exception as e:
                last_error = e
                error_msg = str(e).lower()
                
                # APIクォータ超過エラーまたはレート制限エラーの場合
                quota_errors = ["insufficient_quota", "quota exceeded", "rate_limit"]
                if any(err in error_msg for err in quota_errors):
                    logger.error(f"{self.api_provider} API制限エラー: {e}")
                    # クォータ超過は再試行しても解決しないので、すぐに例外を発生させる
                    raise e
                
                # レート制限エラーの場合は少し待機してから再試行
                if "rate" in error_msg and "limit" in error_msg:
                    logger.warning(f"レート制限エラーが発生しました（試行 {attempt+1}/{self.max_retries+1}）: {e}")
                    import time
                    time.sleep(2 ** attempt)  # 指数バックオフ
                    continue
                
                # その他のエラーの場合
                logger.warning(f"PPTXコード生成中にエラーが発生しました（試行 {attempt+1}/{self.max_retries+1}）: {e}")
                
                # 最後の試行の場合はフォールバックコードを返す
                if attempt == self.max_retries:
                    logger.error("最大試行回数に達しました。フォールバックコードを返します。")
                    return """
# APIエラーが発生したため、基本的なPPTXファイルを生成するコードを返します
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# テンプレートから新しいプレゼンテーションを作成
prs = Presentation('workspace/input/template.pptx')

# タイトルスライドを追加
slide_layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(slide_layout)

# プレースホルダーへテキストを設定
title_placeholder = slide.placeholders[10]
subtitle_placeholder = slide.placeholders[11]
presenter_placeholder = slide.placeholders[12]
title_placeholder.text = "APIエラーが発生しました"
subtitle_placeholder.text = "エラーのため簡易的なプレゼンテーションを生成しました"
presenter_placeholder.text = "AI プレゼンテーション生成"

# 内容スライドを追加
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "プレゼンテーション内容"
tf = body.text_frame
tf.text = "APIエラーが発生したため、本来の内容を生成できませんでした。"

p = tf.add_paragraph()
p.text = "1. 問題が解決するまでこの簡易版プレゼンテーションをご利用ください。"
p.level = 1

p = tf.add_paragraph()
p.text = "2. APIの利用制限を確認してください。"
p.level = 1

# 保存
output_file = 'workspace/output/error_presentation.pptx'
prs.save(output_file)
print(f"プレゼンテーションを {output_file} に保存しました。")
"""
        
        # すべての試行が失敗した場合
        if last_error:
            raise last_error 