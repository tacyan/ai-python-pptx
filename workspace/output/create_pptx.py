from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# テンプレート読み込み
prs = Presentation('workspace/input/template.pptx')
slide_layouts = prs.slide_layouts

# スライド1
slide = prs.slides.add_slide(slide_layouts[2])
title = slide.shapes.title
subtitle = slide.placeholders[11]
title.text = "冒険の始まり！AIエージェントの世界へようこそ！"
subtitle.text = "AIエージェント入門"

img_path = "workspace/input/images/map.jpg" # 適切な画像ファイル名に変更してください
left = Inches(1)
top = Inches(2)
width = Inches(8)
slide.shapes.add_picture(img_path, left, top, width=width)

tf = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(2))
tf.text = "皆さん、こんにちは！今日はAIエージェントという、ワクワクする新しい世界への冒険に一緒に出かけましょう！プログラミングの基礎知識は既にある皆さんなら、きっとこの旅を楽しめるはずです。スマートスピーカーや自動運転など、AIエージェントはすでに私たちの生活に溶け込んでいます。このプレゼンテーションでは、AIエージェントの基本概念から実装方法まで、分かりやすく解説します。"
tf.text_frame.paragraphs[0].font.size = Pt(16)


# スライド2
slide = prs.slides.add_slide(slide_layouts[0])
slide.shapes.title.text = "未知の土地へ！AIエージェントとは？"

left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(3)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
shape.line.fill.fore_color.rgb = RGBColor(0, 0, 0)

# 人間とAIエージェントの相互作用図(簡略化)
# 複雑な図形は外部ツールで作成し、画像として挿入することを推奨
# 以下はテキストによる説明の代用
tf = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(3))
tf.text = "AIエージェントとは、環境を感知し、目標を達成するために自ら行動するプログラムのことです。まるでゲームの主人公のように、自律的に行動します。主な特徴は「知覚」「意思決定」「行動」の3つです。"
tf.text_frame.paragraphs[0].font.size = Pt(16)


# スライド3
slide = prs.slides.add_slide(slide_layouts[0])
slide.shapes.title.text = "迷路を解くAIエージェント"

# 迷路のイラスト(簡略化)  画像挿入推奨
tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
tf.text = "例えば、迷路を解くAIエージェントを考えてみましょう。まず、センサー（知覚）で壁の位置を認識します。次に、最適な経路を計算（意思決定）し、実際に移動（行動）します。このように、AIエージェントは環境を理解し、目標達成のために自律的に行動します。"
tf.text_frame.paragraphs[0].font.size = Pt(16)


# スライド4
slide = prs.slides.add_slide(slide_layouts[0])
slide.shapes.title.text = "分岐点！AIエージェント vs. チャットボット"

rows = 4
cols = 3
left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(3)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table
table.cell(0, 0).text = "項目"
table.cell(1, 0).text = "主な機能"
table.cell(2, 0).text = "目的"
table.cell(3, 0).text = "複雑さ"
table.cell(0, 1).text = "AIエージェント"
table.cell(1, 1).text = "環境感知、意思決定、行動、タスク実行"
table.cell(2, 1).text = "複雑なタスクの自動化"
table.cell(3, 1).text = "高度"
table.cell(0, 2).text = "チャットボット"
table.cell(1, 2).text = "自然言語による会話"
table.cell(2, 2).text = "ユーザーとの対話、情報提供"
table.cell(3, 2).text = "比較的単純"

tf = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
tf.text = "AIエージェントとよく似た言葉に「チャットボット」がありますが、両者は異なります。チャットボットは会話に特化していますが、AIエージェントはもっと幅広いタスクを実行できます。"
tf.text_frame.paragraphs[0].font.size = Pt(16)


# スライド5, 6, 7, 8, 9 は同様の手法で作成


# スライド6: LangChainコード例
slide = prs.slides.add_slide(slide_layouts[0])
slide.shapes.title.text = "魔法のツール！LangChainとLangGraphによる実装"

img_path = "workspace/input/images/langchain_logo.png" # 適切な画像ファイル名に変更してください
left = Inches(1)
top = Inches(2)
width = Inches(4)
slide.shapes.add_picture(img_path, left, top, width=width)

img_path = "workspace/input/images/langgraph_logo.png" # 適切な画像ファイル名に変更してください
left = Inches(5)
top = Inches(2)
width = Inches(4)
slide.shapes.add_picture(img_path, left, top, width=width)

code_text = """
# 簡単なLangChainのコード例（コメント付き）
from langchain.agents import load_tools
from langchain.agents import initialize_agent
from langchain.agents import AgentType
from langchain.llms import OpenAI

# OpenAI APIキーを設定
# ...

# ツールのロード
tools = load_tools(["serpapi", "llm-math"], llm=OpenAI(temperature=0))

# エージェントの初期化
agent = initialize_agent(tools, OpenAI(temperature=0), agent=AgentType.ZERO_SHOT_REACT_DESCRIPTION, verbose=True)

# エージェントの実行
agent.run("東京の人口は？")
"""
tf = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(4))
tf.text = code_text
tf.text_frame.paragraphs[0].font.size = Pt(12)
tf.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


# スライド7以降も同様の手法で作成。画像や図表は適切なファイル名で置き換えてください。

prs.save('workspace/output/presentation.pptx')

