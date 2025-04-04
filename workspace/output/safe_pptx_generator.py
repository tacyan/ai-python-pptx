#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
安全なPowerPointプレゼンテーション生成スクリプト

このスクリプトは、プレースホルダーエラーが発生した場合のフォールバックソリューションとして
PowerPointプレゼンテーションを安全に生成します。テンプレートに依存せず、
テキストボックスと基本的な図形のみを使用して、一意のファイル名でプレゼンテーションを作成します。

制限事項:
- テンプレートファイルを使用しません
- プレースホルダーに依存しません
- 基本的なスライドレイアウトのみを使用します
"""

import os
import sys
import datetime
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def ユニークファイル名生成(プレフィックス="安全_プレゼンテーション", 拡張子="pptx"):
    """
    タイムスタンプとランダムな数値を含む一意のファイル名を生成します。
    
    @param {string} プレフィックス - ファイル名の先頭部分
    @param {string} 拡張子 - ファイルの拡張子
    @return {string} 生成されたファイルパス
    """
    現在時刻 = datetime.datetime.now()
    タイムスタンプ = 現在時刻.strftime("%Y%m%d_%H%M%S")
    ランダム接尾辞 = f"_{random.randint(1000, 9999)}"
    ファイル名 = f"{プレフィックス}_{タイムスタンプ}{ランダム接尾辞}.{拡張子}"
    出力ディレクトリ = os.path.join("workspace", "output")
    os.makedirs(出力ディレクトリ, exist_ok=True)
    return os.path.join(出力ディレクトリ, ファイル名)

def 安全プレゼンテーション作成():
    """
    プレースホルダーに依存しない安全なプレゼンテーションを作成します。
    テキストボックスと基本図形を使用してスライドを構築します。
    
    @return {string} 生成されたプレゼンテーションファイルのパス
    """
    # 新しいプレゼンテーションを作成
    プレゼンテーション = Presentation()
    
    # タイトルスライド
    スライド = プレゼンテーション.slides.add_slide(プレゼンテーション.slide_layouts[0])
    
    # テキストボックスでタイトルを追加
    タイトル図形 = スライド.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    タイトルフレーム = タイトル図形.text_frame
    タイトルフレーム.text = "安全に生成されたプレゼンテーション"
    タイトルフレーム.paragraphs[0].alignment = PP_ALIGN.CENTER
    タイトルフレーム.paragraphs[0].font.size = Pt(44)
    タイトルフレーム.paragraphs[0].font.bold = True
    タイトルフレーム.paragraphs[0].font.color.rgb = RGBColor(0, 75, 120)
    
    # サブタイトル
    サブタイトル図形 = スライド.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    サブタイトルフレーム = サブタイトル図形.text_frame
    サブタイトルフレーム.text = "プレースホルダーエラーを回避する安全な実装"
    サブタイトルフレーム.paragraphs[0].alignment = PP_ALIGN.CENTER
    サブタイトルフレーム.paragraphs[0].font.size = Pt(28)
    サブタイトルフレーム.paragraphs[0].font.color.rgb = RGBColor(70, 70, 70)
    
    # 日付情報
    日付図形 = スライド.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(0.5))
    日付フレーム = 日付図形.text_frame
    現在日時 = datetime.datetime.now().strftime("%Y年%m月%d日 %H:%M:%S")
    日付フレーム.text = f"作成日時: {現在日時}"
    日付フレーム.paragraphs[0].alignment = PP_ALIGN.CENTER
    日付フレーム.paragraphs[0].font.size = Pt(14)
    日付フレーム.paragraphs[0].font.italic = True
    
    # 背景を彩るための図形を追加
    左上図形 = スライド.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, 
        Inches(0), Inches(0), 
        Inches(2), Inches(2)
    )
    左上図形.fill.solid()
    左上図形.fill.fore_color.rgb = RGBColor(230, 240, 250)
    左上図形.line.color.rgb = RGBColor(200, 220, 240)
    
    右下図形 = スライド.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(8), Inches(5), 
        Inches(2), Inches(1)
    )
    右下図形.fill.solid()
    右下図形.fill.fore_color.rgb = RGBColor(240, 240, 220)
    右下図形.line.color.rgb = RGBColor(220, 220, 200)
    
    # コンテンツスライド
    スライド = プレゼンテーション.slides.add_slide(プレゼンテーション.slide_layouts[1])
    
    # タイトル
    タイトル図形 = スライド.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    タイトルフレーム = タイトル図形.text_frame
    タイトルフレーム.text = "安全なプレゼンテーション生成"
    タイトルフレーム.paragraphs[0].alignment = PP_ALIGN.CENTER
    タイトルフレーム.paragraphs[0].font.size = Pt(40)
    タイトルフレーム.paragraphs[0].font.bold = True
    タイトルフレーム.paragraphs[0].font.color.rgb = RGBColor(0, 75, 120)
    
    # 背景図形
    背景図形 = スライド.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.5), Inches(1.8), 
        Inches(9), Inches(4.2)
    )
    背景図形.fill.solid()
    背景図形.fill.fore_color.rgb = RGBColor(245, 245, 250)
    背景図形.line.color.rgb = RGBColor(220, 220, 230)
    
    # コンテンツ
    コンテンツ図形 = スライド.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    コンテンツフレーム = コンテンツ図形.text_frame
    
    段落 = コンテンツフレーム.paragraphs[0]
    段落.text = "• プレースホルダーエラーを回避するために安全モードで生成"
    段落.font.size = Pt(24)
    段落.font.color.rgb = RGBColor(50, 50, 50)
    
    段落 = コンテンツフレーム.add_paragraph()
    段落.text = "• テキストボックスと図形のみを使用"
    段落.font.size = Pt(24)
    段落.font.color.rgb = RGBColor(50, 50, 50)
    
    段落 = コンテンツフレーム.add_paragraph()
    段落.text = "• 日付とタイムスタンプを含む一意なファイル名"
    段落.font.size = Pt(24)
    段落.font.color.rgb = RGBColor(50, 50, 50)
    
    段落 = コンテンツフレーム.add_paragraph()
    段落.text = "• エラー処理を強化し安定した動作を確保"
    段落.font.size = Pt(24)
    段落.font.color.rgb = RGBColor(50, 50, 50)
    
    # 図解スライド
    スライド = プレゼンテーション.slides.add_slide(プレゼンテーション.slide_layouts[1])
    
    # タイトル
    タイトル図形 = スライド.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    タイトルフレーム = タイトル図形.text_frame
    タイトルフレーム.text = "安全実装の図解"
    タイトルフレーム.paragraphs[0].alignment = PP_ALIGN.CENTER
    タイトルフレーム.paragraphs[0].font.size = Pt(40)
    タイトルフレーム.paragraphs[0].font.bold = True
    タイトルフレーム.paragraphs[0].font.color.rgb = RGBColor(0, 75, 120)
    
    # 図解の配置
    # 中央の円
    中央円 = スライド.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(4), Inches(2.5), 
        Inches(2), Inches(2)
    )
    中央円.fill.solid()
    中央円.fill.fore_color.rgb = RGBColor(0, 112, 192)
    中央円.line.color.rgb = RGBColor(0, 80, 160)
    
    # 中央円のテキスト
    中央テキスト = スライド.shapes.add_textbox(
        Inches(4.25), Inches(3.2), 
        Inches(1.5), Inches(0.6)
    )
    中央テキストフレーム = 中央テキスト.text_frame
    中央テキストフレーム.text = "安全な\n生成方式"
    中央テキストフレーム.paragraphs[0].alignment = PP_ALIGN.CENTER
    中央テキストフレーム.paragraphs[0].font.size = Pt(14)
    中央テキストフレーム.paragraphs[0].font.bold = True
    中央テキストフレーム.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # 左の要素
    左図形 = スライド.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(1), Inches(2.5), 
        Inches(2), Inches(1)
    )
    左図形.fill.solid()
    左図形.fill.fore_color.rgb = RGBColor(240, 240, 240)
    左図形.line.color.rgb = RGBColor(200, 200, 200)
    
    左テキスト = スライド.shapes.add_textbox(
        Inches(1.1), Inches(2.7), 
        Inches(1.8), Inches(0.6)
    )
    左テキストフレーム = 左テキスト.text_frame
    左テキストフレーム.text = "テキストボックス\n直接配置"
    左テキストフレーム.paragraphs[0].alignment = PP_ALIGN.CENTER
    左テキストフレーム.paragraphs[0].font.size = Pt(12)
    左テキストフレーム.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
    
    # 右の要素
    右図形 = スライド.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(7), Inches(2.5), 
        Inches(2), Inches(1)
    )
    右図形.fill.solid()
    右図形.fill.fore_color.rgb = RGBColor(240, 240, 240)
    右図形.line.color.rgb = RGBColor(200, 200, 200)
    
    右テキスト = スライド.shapes.add_textbox(
        Inches(7.1), Inches(2.7), 
        Inches(1.8), Inches(0.6)
    )
    右テキストフレーム = 右テキスト.text_frame
    右テキストフレーム.text = "基本図形のみ\n使用"
    右テキストフレーム.paragraphs[0].alignment = PP_ALIGN.CENTER
    右テキストフレーム.paragraphs[0].font.size = Pt(12)
    右テキストフレーム.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
    
    # 上の要素
    上図形 = スライド.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(4), Inches(1), 
        Inches(2), Inches(1)
    )
    上図形.fill.solid()
    上図形.fill.fore_color.rgb = RGBColor(240, 240, 240)
    上図形.line.color.rgb = RGBColor(200, 200, 200)
    
    上テキスト = スライド.shapes.add_textbox(
        Inches(4.1), Inches(1.2), 
        Inches(1.8), Inches(0.6)
    )
    上テキストフレーム = 上テキスト.text_frame
    上テキストフレーム.text = "プレースホルダー\n不使用"
    上テキストフレーム.paragraphs[0].alignment = PP_ALIGN.CENTER
    上テキストフレーム.paragraphs[0].font.size = Pt(12)
    上テキストフレーム.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
    
    # 下の要素
    下図形 = スライド.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(4), Inches(5), 
        Inches(2), Inches(1)
    )
    下図形.fill.solid()
    下図形.fill.fore_color.rgb = RGBColor(240, 240, 240)
    下図形.line.color.rgb = RGBColor(200, 200, 200)
    
    下テキスト = スライド.shapes.add_textbox(
        Inches(4.1), Inches(5.2), 
        Inches(1.8), Inches(0.6)
    )
    下テキストフレーム = 下テキスト.text_frame
    下テキストフレーム.text = "ユニークな\nファイル名生成"
    下テキストフレーム.paragraphs[0].alignment = PP_ALIGN.CENTER
    下テキストフレーム.paragraphs[0].font.size = Pt(12)
    下テキストフレーム.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
    
    # 線を引いて接続
    左線 = スライド.shapes.add_connector(
        1, Inches(3), Inches(3), Inches(4), Inches(3.5)
    )
    左線.line.color.rgb = RGBColor(100, 100, 100)
    
    右線 = スライド.shapes.add_connector(
        1, Inches(7), Inches(3), Inches(6), Inches(3.5)
    )
    右線.line.color.rgb = RGBColor(100, 100, 100)
    
    上線 = スライド.shapes.add_connector(
        1, Inches(5), Inches(2), Inches(5), Inches(2.5)
    )
    上線.line.color.rgb = RGBColor(100, 100, 100)
    
    下線 = スライド.shapes.add_connector(
        1, Inches(5), Inches(5), Inches(5), Inches(4.5)
    )
    下線.line.color.rgb = RGBColor(100, 100, 100)
    
    # 保存
    出力ファイル名 = ユニークファイル名生成()
    try:
        プレゼンテーション.save(出力ファイル名)
        print(f"プレゼンテーションが正常に生成されました！保存先: {出力ファイル名}")
        return 出力ファイル名
    except Exception as エラー:
        print(f"保存中にエラーが発生しました: {エラー}")
        # フォールバックとして別の名前で保存を試みる
        フォールバック名 = os.path.join("workspace", "output", f"安全プレゼンテーション_緊急_{random.randint(1000, 9999)}.pptx")
        try:
            プレゼンテーション.save(フォールバック名)
            print(f"フォールバックで保存に成功しました: {フォールバック名}")
            return フォールバック名
        except Exception as 二次エラー:
            print(f"フォールバック保存にも失敗しました: {二次エラー}")
            return None

if __name__ == "__main__":
    try:
        出力ファイル = 安全プレゼンテーション作成()
        if 出力ファイル:
            print(f"ファイルパス: {出力ファイル}")
            sys.exit(0)
        else:
            print("プレゼンテーションの生成に失敗しました")
            sys.exit(1)
    except Exception as エラー:
        print(f"エラーが発生しました: {エラー}")
        sys.exit(1)
