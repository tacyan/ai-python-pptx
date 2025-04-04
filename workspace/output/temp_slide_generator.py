#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
PowerPointプレゼンテーション自動生成スクリプト
このスクリプトは完全に新規にPowerPointプレゼンテーションを作成します。
テンプレートファイルを使用せず、プレースホルダーにアクセスせず、
すべてのコンテンツをテキストボックスと図形で直接作成します。
タイムスタンプ付きファイル名で常に新規ファイルを生成します。

特徴:
- テンプレート非依存の完全新規作成
- プレースホルダーを一切使用しない
- 堅牢なエラーハンドリング
- 一意なタイムスタンプとランダム値のファイル名
- workspace/output/create_pptx.pyを変更せずに動作

Author: AIアシスタント
Version: 1.0
Date: 2025-04-04
"""

import os
import sys
import datetime
import random
import traceback

# 必要なライブラリを確認してインポート
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptxライブラリがインストールされていません。")
    print("次のコマンドでインストールしてください: pip install python-pptx")
    sys.exit(1)

# 一意のファイル名を生成
def generate_unique_filename(prefix="Presentation", ext="pptx"):
    """
    タイムスタンプとランダムな値を使用して一意のファイル名を生成します
    
    Args:
        prefix (str): ファイル名の接頭辞
        ext (str): ファイルの拡張子
        
    Returns:
        str: 生成されたファイルパス
    """
    now = datetime.datetime.now()
    timestamp = now.strftime("%Y%m%d_%H%M%S")
    random_suffix = f"_{random.randint(1000, 9999)}"
    filename = f"{prefix}_{timestamp}{random_suffix}.{ext}"
    
    # 出力ディレクトリの確認と作成
    output_dir = os.path.join("workspace", "output")
    os.makedirs(output_dir, exist_ok=True)
    
    return os.path.join(output_dir, filename)

def create_presentation():
    """
    新しいプレゼンテーションを作成します
    
    Returns:
        str: 作成されたプレゼンテーションのファイルパス
    """
    try:
        # 新しいプレゼンテーションを作成
        prs = Presentation()
        
        # プライマリカラーの定義
        primary_color = RGBColor(0, 112, 192)  # 青
        secondary_color = RGBColor(255, 192, 0)  # 黄色
        accent_color = RGBColor(112, 48, 160)  # 紫
        
        # 1. タイトルスライド
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # 標準の最初のレイアウト
        
        # テキストボックスでタイトルを追加 (プレースホルダーを使用しない)
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
        title_frame = title_shape.text_frame
        title_frame.text = "自動生成プレゼンテーション"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = primary_color
        
        # サブタイトル
        subtitle_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.text = "安全なプレースホルダー非依存アプローチ"
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_frame.paragraphs[0].font.size = Pt(32)
        subtitle_frame.paragraphs[0].font.italic = True
        
        # 作成日時
        current_date = datetime.datetime.now().strftime("%Y年%m月%d日 %H:%M:%S")
        date_shape = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(0.5))
        date_frame = date_shape.text_frame
        date_frame.text = f"作成日時: {current_date}"
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        date_frame.paragraphs[0].font.size = Pt(16)
        
        # 2. 特徴スライド
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 標準の2番目のレイアウト
        
        # タイトル
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "このプレゼンテーションの特徴"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        
        # コンテンツ
        content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        content_frame = content_shape.text_frame
        content_frame.word_wrap = True
        
        features = [
            "• テンプレートファイルに依存しない完全新規作成",
            "• プレースホルダーを使用しないテキストボックスアプローチ",
            "• タイムスタンプと乱数を用いた一意なファイル名",
            "• 既存ファイルを上書きしない安全な設計",
            "• 堅牢なエラーハンドリング機能"
        ]
        
        # 最初のパラグラフ
        p = content_frame.paragraphs[0]
        p.text = features[0]
        p.font.size = Pt(24)
        
        # 残りの機能を追加
        for feature in features[1:]:
            p = content_frame.add_paragraph()
            p.text = feature
            p.font.size = Pt(24)
            p.level = 0
        
        # 3. 図形スライド
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # タイトル
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "図形とテキストボックスの例"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        
        # 円形を追加
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(2), Inches(2), Inches(2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = primary_color
        circle.line.color.rgb = RGBColor(0, 0, 0)
        
        # 四角形を追加
        rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(2), Inches(2), Inches(2))
        rectangle.fill.solid()
        rectangle.fill.fore_color.rgb = secondary_color
        rectangle.line.color.rgb = RGBColor(0, 0, 0)
        
        # 三角形を追加（MSO_SHAPE.TRIANGLEが存在しないため、RIGHT_TRIANGLEを使用）
        triangle = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(3.5), Inches(4), Inches(2), Inches(2))
        triangle.fill.solid()
        triangle.fill.fore_color.rgb = accent_color
        triangle.line.color.rgb = RGBColor(0, 0, 0)
        
        # 説明テキストを追加
        text_shape = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
        text_frame = text_shape.text_frame
        text_frame.text = "図形とテキストボックスを使用したスライド作成例"
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.paragraphs[0].font.size = Pt(20)
        text_frame.paragraphs[0].font.italic = True
        
        # 4. まとめスライド
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # タイトル
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "まとめ"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        
        # コンテンツ
        content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        content_frame = content_shape.text_frame
        
        p = content_frame.paragraphs[0]
        p.text = "このプレゼンテーションは以下の手法で作成されました："
        p.font.size = Pt(24)
        
        bullet_points = [
            "• python-pptxライブラリを使用した直接生成",
            "• テンプレートファイルやプレースホルダーに依存しない",
            "• テキストボックスと図形を直接配置",
            "• 毎回新しいファイル名で保存"
        ]
        
        for point in bullet_points:
            p = content_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(24)
            p.level = 0
        
        # フッター情報
        footer_shape = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        footer_frame = footer_shape.text_frame
        footer_frame.text = "© 2025 AI自動生成プレゼンテーション"
        footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        footer_frame.paragraphs[0].font.size = Pt(12)
        
        # プレゼンテーションを保存
        output_filename = generate_unique_filename("Safe_Presentation")
        prs.save(output_filename)
        print(f"プレゼンテーションが正常に生成されました！保存先: {output_filename}")
        return output_filename
        
    except Exception as e:
        print(f"エラーが発生しました: {e}")
        traceback.print_exc()
        
        try:
            # エラー時の代替プレゼンテーション
            print("代替プレゼンテーションを作成します...")
            error_ppt_filename = generate_unique_filename("Error_Presentation")
            
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            
            # エラータイトル
            title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
            title_frame = title_shape.text_frame
            title_frame.text = "エラーが発生しました"
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title_frame.paragraphs[0].font.size = Pt(44)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)  # 赤色
            
            # エラーメッセージ
            error_shape = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
            error_frame = error_shape.text_frame
            error_frame.text = f"プレゼンテーション生成中にエラーが発生しました:\n{str(e)}"
            error_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            error_frame.paragraphs[0].font.size = Pt(20)
            
            # タイムスタンプ
            now = datetime.datetime.now().strftime("%Y年%m月%d日 %H:%M:%S")
            date_shape = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
            date_frame = date_shape.text_frame
            date_frame.text = f"エラー発生時刻: {now}"
            date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            date_frame.paragraphs[0].font.size = Pt(14)
            
            prs.save(error_ppt_filename)
            print(f"エラー用プレゼンテーションが生成されました: {error_ppt_filename}")
            return error_ppt_filename
            
        except Exception as err:
            print(f"代替プレゼンテーションの作成中にエラーが発生しました: {err}")
            traceback.print_exc()
            return None

if __name__ == "__main__":
    try:
        output_file = create_presentation()
        if output_file:
            print(f"プレゼンテーションの生成が完了しました！")
            print(f"ファイルパス: {output_file}")
            sys.exit(0)
        else:
            print("プレゼンテーションの生成に失敗しました。")
            sys.exit(1)
    except Exception as e:
        print(f"予期せぬエラーが発生しました: {e}")
        traceback.print_exc()
        sys.exit(1) 