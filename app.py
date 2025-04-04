"""
このモジュールはパワーポイント資料を自動生成するStreamlitアプリケーションです。
ユーザーがアップロードしたテキストファイルからプレゼンテーション内容を読み取り、
AIエージェントを使用して、PowerPointスライドを生成するPythonコードを出力します。
画像や動画の挿入にも対応しています。
OpenAIとGemini両方のAPIに対応しています。

Usage:
    streamlit run app.py
"""

import os
import streamlit as st
import sys
import subprocess
import shutil
import importlib
import time
import logging
import requests
import re

# Streamlitページ設定を最初に行う（このアプリ全体で一度だけ）
st.set_page_config(page_title="AIプレゼンテーション生成", page_icon="📊", layout="wide")

# ロギングの設定
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# 必要なライブラリが存在するか確認し、なければインストールする
def check_and_install_dependencies():
    """
    必要な依存関係をチェックし、不足している場合はインストールする
    エラーハンドリングを強化し、インストール失敗時の代替手段も提供
    """
    # インストール履歴ファイルの確認
    install_history_file = "installed_packages.txt"
    installed_packages = set()
    
    if os.path.exists(install_history_file):
        try:
            with open(install_history_file, 'r') as f:
                installed_packages = set(line.strip() for line in f if line.strip())
            
            st.info("前回のインストール情報を読み込みました。必要なパッケージのみインストールします。")
        except Exception as e:
            st.warning(f"インストール履歴の読み込みに失敗しました。すべてのパッケージを確認します: {e}")
    
    required_packages = [
        "langchain-core==0.3.0",
        "langchain-openai==0.2.0",
        "langchain-community>=0.0.10",
        "langgraph==0.2.22",
        "python-pptx==1.0.2",
        "ipython"
    ]
    
    # Gemini関連のパッケージ（オプション）
    gemini_packages = [
        "google-generativeai>=0.3.0",
        "langchain-google-genai==0.1.5"
    ]
    
    # 未インストールのパッケージのみを選択
    to_install = [pkg for pkg in required_packages if pkg.split('==')[0].split('>=')[0] not in installed_packages]
    
    if not to_install:
        st.success("必要なパッケージはすべてインストール済みです。")
    else:
        st.info(f"{len(to_install)}個のパッケージをインストールします...")
    
    # 基本パッケージのインストール
    all_success = install_packages(to_install, critical=True, install_history=installed_packages)
    
    # 特別な方法でGeminiパッケージをインストール
    gemini_available = install_gemini_packages(installed_packages)
    
    # インストール履歴を保存
    try:
        with open(install_history_file, 'w') as f:
            for package in installed_packages:
                f.write(f"{package}\n")
    except Exception as e:
        st.warning(f"インストール履歴の保存に失敗しました: {e}")
    
    return gemini_available

def install_gemini_packages(install_history=None):
    """
    Gemini API関連パッケージを特別な方法でインストールする
    
    Args:
        install_history (set): インストール済みパッケージの集合
    
    Returns:
        bool: Geminiパッケージのインストールに成功したかどうか
    """
    if install_history is None:
        install_history = set()
    
    # すでにインストールされているか確認
    if "google-generativeai" in install_history and "langchain-google-genai" in install_history:
        st.success("Google Gemini APIのサポートは既に有効化されています！")
        return True
    
    try:
        # まずgoogle-generativeaiをインストール
        if "google-generativeai" not in install_history:
            st.info("Google AI SDKをインストールしています...")
            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "--no-cache-dir", "google-generativeai>=0.3.0"],
                stderr=subprocess.PIPE
            )
            install_history.add("google-generativeai")
        
        # 依存関係のチェック
        if "langchain-google-genai" not in install_history:
            st.info("LangChain Google Genai拡張をインストールしています...")
            
            # 複数の方法を試す
            install_success = False
            
            methods = [
                # 方法1: 最新バージョンで直接インストール
                [sys.executable, "-m", "pip", "install", "--no-cache-dir", "langchain-google-genai"],
                # 方法2: 特定バージョンを指定してインストール
                [sys.executable, "-m", "pip", "install", "--no-cache-dir", "--use-pep517", "langchain-google-genai==0.1.5"],
                # 方法3: 依存関係を明示的に指定
                [sys.executable, "-m", "pip", "install", "--no-cache-dir", "langchain-google-genai==0.1.5", "typing-inspect>=0.8.0", "typing-extensions>=4.5.0"],
                # 方法4: 最小限の依存関係でインストール
                [sys.executable, "-m", "pip", "install", "--no-cache-dir", "--no-deps", "langchain-google-genai", "google-ai-generativelanguage>=0.6.0"]
            ]
            
            for i, cmd in enumerate(methods):
                try:
                    st.info(f"インストール方法 {i+1}/{len(methods)} を試行中...")
                    subprocess.check_call(cmd, stderr=subprocess.PIPE)
                    install_success = True
                    st.success(f"インストール方法 {i+1} で成功しました！")
                    install_history.add("langchain-google-genai")
                    break
                except subprocess.CalledProcessError as e:
                    st.warning(f"インストール方法 {i+1} が失敗しました: {e}")
                    continue
            
            # インポートテスト
            try:
                import langchain_google_genai
                st.success("Google Gemini APIのサポートが有効化されました！")
                return True
            except ImportError as e:
                if "langchain-google-genai" in install_history:
                    st.warning(f"パッケージはインストールされましたが、インポートできませんでした: {e}")
                    st.info("アプリケーションを再起動して再試行してください。")
                else:
                    st.error("すべてのインストール方法が失敗しました。")
                return False
            
    except Exception as e:
        st.warning(f"Google Gemini APIサポートのインストールに失敗しました: {e}")
        st.info("OpenAI APIのみ使用可能です。手動インストールを試すには次のコマンドを実行してください:")
        st.code("pip install langchain-google-genai --no-cache-dir")
        return False

def install_packages(packages, critical=False, install_history=None):
    """
    パッケージのリストをインストールする

    Args:
        packages (list): インストールするパッケージのリスト
        critical (bool): 重要なパッケージかどうか。Trueの場合、インストール失敗時にエラーを表示
        install_history (set): インストール済みパッケージの集合（Noneの場合は作成）

    Returns:
        bool: すべてのパッケージのインストールに成功したかどうか
    """
    if install_history is None:
        install_history = set()
    
    all_success = True
    
    for package in packages:
        try:
            package_name = package.split("==")[0].split(">=")[0]
            
            # すでに履歴にあればスキップ
            if package_name in install_history:
                logger.info(f"{package_name} は既にインストールされています")
                continue
                
            try:
                # すでにインストールされているか確認
                importlib.import_module(package_name.replace("-", "_"))
                logger.info(f"{package_name} は既にインストールされています")
                install_history.add(package_name)
                continue
            except ImportError:
                st.info(f"{package_name} をインストールしています...")
                
                # インストール試行回数
                max_attempts = 3
                
                for attempt in range(max_attempts):
                    try:
                        # インストールコマンドを実行 - オプションを追加して成功率を高める
                        subprocess.check_call(
                            [sys.executable, "-m", "pip", "install", "--no-cache-dir", "--upgrade", package],
                            stderr=subprocess.STDOUT
                        )
                        st.success(f"{package_name} のインストールが完了しました！")
                        install_history.add(package_name)
                        break
                    except subprocess.CalledProcessError as e:
                        logger.warning(f"{package_name} のインストールに失敗しました（試行 {attempt+1}/{max_attempts}）: {e}")
                        
                        if attempt == max_attempts - 1:
                            # すべての試行が失敗
                            if critical:
                                st.error(f"重要なパッケージ {package_name} のインストールに失敗しました。")
                                st.error(f"手動でインストールするには: pip install {package}")
                                all_success = False
                            else:
                                st.warning(f"オプションパッケージ {package_name} のインストールに失敗しました。一部の機能が制限される場合があります。")
                                all_success = False
                        else:
                            # 短い遅延後に再試行
                            time.sleep(2)
        except Exception as e:
            logger.error(f"パッケージ {package} の処理中にエラーが発生しました: {e}")
            if critical:
                st.error(f"パッケージ {package} の処理中にエラーが発生しました: {e}")
                all_success = False
            else:
                st.warning(f"オプションパッケージ {package} の処理中にエラーが発生しました。一部の機能が制限される場合があります。")
                all_success = False
    
    return all_success

# ディレクトリ構造を確認し、存在しない場合は作成
def ensure_directories():
    """
    必要なディレクトリ構造を確保する
    """
    os.makedirs("workspace/input", exist_ok=True)
    os.makedirs("workspace/output", exist_ok=True)
    os.makedirs("workspace/input/images", exist_ok=True)

# 安全なプレゼンテーション生成関数をここに直接定義（サブプロセスで呼び出す代わりに）
def generate_safe_presentation():
    """
    プレースホルダーエラーを回避するシンプルなプレゼンテーションを生成する安全な関数
    
    Returns:
        str: 生成されたファイルのパス
    """
    import os
    import datetime
    import random
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    # 一意のファイル名を生成
    now = datetime.datetime.now()
    timestamp = now.strftime("%Y%m%d_%H%M%S")
    random_suffix = f"_{random.randint(1000, 9999)}"
    filename = f"Safe_Presentation_{timestamp}{random_suffix}.pptx"
    output_dir = os.path.join("workspace", "output")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, filename)
    
    # 新しいプレゼンテーションを作成
    prs = Presentation()
    
    # タイトルスライド
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # テキストボックスでタイトルを追加
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    title_frame = title_shape.text_frame
    title_frame.text = "安全に生成されたプレゼンテーション"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    
    # サブタイトル
    subtitle_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "プレースホルダーエラーを回避する安全な実装"
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    
    # 日付情報
    date_shape = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(0.5))
    date_frame = date_shape.text_frame
    current_date = datetime.datetime.now().strftime("%Y年%m月%d日 %H:%M:%S")
    date_frame.text = f"作成日時: {current_date}"
    date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # コンテンツスライド
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # タイトル
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "安全なプレゼンテーション生成"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    
    # コンテンツ
    content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    content_frame = content_shape.text_frame
    
    p = content_frame.paragraphs[0]
    p.text = "• プレースホルダーエラーを回避するために安全モードで生成"
    p.font.size = Pt(24)
    
    p = content_frame.add_paragraph()
    p.text = "• テキストボックスと図形のみを使用"
    p.font.size = Pt(24)
    
    p = content_frame.add_paragraph()
    p.text = "• 日付とタイムスタンプを含む一意なファイル名"
    p.font.size = Pt(24)
    
    # 保存
    prs.save(output_path)
    return output_path

def main():
    """
    Streamlitアプリのメイン関数
    """
    # 依存関係のチェックとインストール
    gemini_available = check_and_install_dependencies()
    
    # モジュールのインポート
    try:
        from langchain_openai import ChatOpenAI
        if gemini_available:
            try:
                from langchain_google_genai import ChatGoogleGenerativeAI
            except ImportError:
                gemini_available = False
        from pptx_agent import PPTXAgent
    except ImportError as e:
        st.error(f"必要なライブラリをインポートできませんでした: {e}")
        st.error("アプリケーションを再起動してください。")
        return
    
    # ディレクトリの確認
    ensure_directories()
    
    # タイトルと説明
    st.title("AIプレゼンテーション自動生成アプリ")
    st.write("プレゼンテーションの内容を記述したテキストファイルをアップロードすると、AIがパワーポイントの資料を自動生成します.")
    
    # ファイルアップロード
    uploaded_file = st.file_uploader("テキストファイルをアップロード (.txt, .docx, .md)", type=["txt", "docx", "md"])
    
    # APIプロバイダー選択
    api_provider_options = ["OpenAI"]
    if gemini_available:
        api_provider_options.append("Google Gemini")
    else:
        st.warning("Google Gemini APIの依存関係をインストールできませんでした。OpenAI APIのみ使用可能です。")
    
    api_provider = st.radio("使用するAPIプロバイダー", api_provider_options, horizontal=True)
    
    # APIキーの入力と関連設定
    if api_provider == "OpenAI":
        api_key = st.text_input("OpenAI APIキー", type="password")
        model = st.selectbox("使用するモデル", options=["gpt-4o", "gpt-3.5-turbo"], index=0)
    else:  # Gemini
        api_key = st.text_input("Google Gemini APIキー", type="password")
        model = st.selectbox("使用するモデル", options=["gemini-pro", "gemini-1.5-pro", "gemini-1.5-flash"], index=0)
    
    # バックアップオプション
    use_fallback = st.checkbox("APIクォータ超過時にバックアップモデルを使用する", value=True, 
                             help="APIクォータが超過した場合、より軽量なモデルにフォールバックします")
    
    # テンプレートファイルの情報
    st.info("注意: このアプリケーションは `workspace/input/template.pptx` というパワーポイントテンプレートファイルを必要とします.")
    
    # テンプレートファイルのアップロード
    template_file = st.file_uploader("テンプレートファイルをアップロード (.pptx)", type=["pptx"])
    if template_file:
        with open("workspace/input/template.pptx", "wb") as f:
            f.write(template_file.getbuffer())
        st.success("テンプレートファイルがアップロードされました！")
    
    # 画像・動画ファイルのアップロード（複数可）
    st.subheader("画像・動画ファイルのアップロード(オプション)")
    st.write("プレゼンテーションに使用する画像や動画ファイルをアップロードできます.")
    
    # 既存の画像・動画ファイルを表示
    existing_media_files = os.listdir("workspace/input/images")
    if existing_media_files:
        st.write("現在アップロードされているファイル:")
        cols = st.columns(4)
        for i, file in enumerate(existing_media_files):
            file_path = f"workspace/input/images/{file}"
            file_type = file.split(".")[-1].lower()
            with cols[i % 4]:
                if file_type in ["jpg", "jpeg", "png", "gif"]:
                    st.image(file_path, caption=file, width=150)
                elif file_type in ["mp4", "mov", "avi"]:
                    st.video(file_path)
                else:
                    st.write(f"📁 {file}")
                
                # ファイル削除ボタン
                if st.button(f"削除: {file}", key=f"delete_{file}"):
                    os.remove(file_path)
                    st.success(f"{file} を削除しました。")
                    st.experimental_rerun()
    
    # 新しいファイルのアップロード
    uploaded_media_files = st.file_uploader("画像・動画ファイルをアップロード (.jpg, .jpeg, .png, .gif, .mp4, .mov)", 
                                          type=["jpg", "jpeg", "png", "gif", "mp4", "mov"], 
                                          accept_multiple_files=True)
    
    if uploaded_media_files:
        for uploaded_file in uploaded_media_files:
            file_path = f"workspace/input/images/{uploaded_file.name}"
            with open(file_path, "wb") as f:
                f.write(uploaded_file.getbuffer())
        
        st.success(f"{len(uploaded_media_files)}個のファイルがアップロードされました！")
        st.experimental_rerun()
    
    # main関数の外側のtryブロック開始
    try:
        # 実行ボタン
        if st.button("プレゼンテーションを生成", key="generate"):
            # テンプレートファイルのチェック
            if not os.path.exists("workspace/input/template.pptx"):
                st.warning("テンプレートファイルが見つかりません。新規プレゼンテーションとして生成します。")
            
            with st.spinner("AIがプレゼンテーションを生成中..."):
                # AI生成とファイル処理のtryブロック開始
                try:
                    # APIキーの設定
                    if api_provider == "OpenAI":
                        os.environ["OPENAI_API_KEY"] = api_key
                    else:  # Gemini
                        os.environ["GOOGLE_API_KEY"] = api_key
                    
                    # アップロードされたファイルを保存
                    content = uploaded_file.read().decode("utf-8")
                    
                    # LLMモデルを初期化
                    try:
                        if api_provider == "OpenAI":
                            llm = ChatOpenAI(model=model, temperature=0.0)
                            fallback_model = "gpt-3.5-turbo" if model != "gpt-3.5-turbo" else None
                        else:  # Gemini
                            llm = ChatGoogleGenerativeAI(model=model, temperature=0.0)
                            fallback_model = "gemini-1.5-flash" if model != "gemini-1.5-flash" else None
                        
                        # PPTXAgentを初期化
                        agent = PPTXAgent(llm=llm, use_fallback=use_fallback, api_provider=api_provider, fallback_model=fallback_model)
                        
                        # エージェントを実行して最終的な出力を取得
                        final_output = agent.run(user_request=content)
                    except Exception as api_error:
                        # APIクォータ超過エラー・レート制限エラーの処理
                        error_msg = str(api_error).lower()
                        quota_error = any(err in error_msg for err in ["insufficient_quota", "rate_limit", "quota exceeded"])
                        
                        if quota_error:
                            st.warning(f"{api_provider} APIの制限に達しました。バックアップモデルを使用します。")
                            if use_fallback:
                                if api_provider == "OpenAI" and model != "gpt-3.5-turbo":
                                    llm = ChatOpenAI(model="gpt-3.5-turbo", temperature=0.0)
                                    agent = PPTXAgent(llm=llm, use_fallback=use_fallback, api_provider=api_provider)
                                    final_output = agent.run(user_request=content)
                                elif api_provider == "Google Gemini" and model != "gemini-1.5-flash":
                                    llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash", temperature=0.0)
                                    agent = PPTXAgent(llm=llm, use_fallback=use_fallback, api_provider=api_provider)
                                    final_output = agent.run(user_request=content)
                                else:
                                    st.error(f"{api_provider} APIの制限に達しました。APIキーの制限を確認するか、別のAPIキーを使用してください。")
                                    st.error("この問題を解決するには：")
                                    st.error(f"1. {api_provider}ダッシュボードでAPIの使用状況と制限を確認する")
                                    st.error("2. 有料プランにアップグレードする")
                                    st.error("3. 別のAPIキーを使用する")
                                    return
                            else:
                                st.error(f"{api_provider} APIの制限に達しました。APIキーの制限を確認するか、別のAPIキーを使用してください。")
                                st.error("この問題を解決するには：")
                                st.error(f"1. {api_provider}ダッシュボードでAPIの使用状況と制限を確認する")
                                st.error("2. 有料プランにアップグレードする")
                                st.error("3. 別のAPIキーを使用する")
                                return
                        else:
                            raise api_error
                    
                    # Python コードブロックが含まれている場合の処理
                    if "```python" in final_output:
                        final_output = final_output.split("```python\n")[-1].split("```")[0]
                        
                        # 出力をファイルに保存（UTF-8エンコーディングを明示的に指定）
                        with open("workspace/output/create_pptx.py", "w", encoding="utf-8") as f:
                            f.write(final_output)
                        
                        # Pythonコードを表示
                        with st.expander("生成されたPythonコード"):
                            st.code(final_output, language="python")
                        
                        # 変数初期化 - 常に modified_code を設定
                        modified_code = final_output  # デフォルトでは最初のコード生成物をそのまま使用
                        
                        # プレースホルダーエラーに対処するコードを追加
                        # 出力コード内にプレースホルダーへの参照があれば、安全なアクセス方法に置換
                        if "placeholder" in final_output and "slide.placeholders" in final_output:
                            st.info("プレースホルダーを使用するコードが検出されました。安全なアクセス方法に変換します。")
                            # 安全なプレースホルダーアクセスコードを追加する部分
                            safe_placeholder_code = """
# 必要なインポート
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 塗りつぶしに関する安全なアクセスのためのヘルパー関数
def set_fill_color_safe(fill, color):
    '''フィルに安全に色を設定する(Noneフィルの場合はsolid()を先に呼び出す)'''
    try:
        if hasattr(fill, 'type') and fill.type == None:
            fill.solid()
        fill.fore_color.rgb = color
        return True
    except (AttributeError, TypeError) as e:
        print(f"フィルの色設定に失敗しました: {e}")
        try:
            # 別の方法を試す
            fill.solid()
            fill.fore_color.rgb = color
            return True
        except Exception as e2:
            print(f"フィルの色設定の2回目の試行も失敗しました: {e2}")
            return False

# 画像ファイルを安全に扱うためのヘルパー関数
def add_image_safe(slide, image_path, left=Inches(1), top=Inches(2), width=Inches(4), height=Inches(3)):
    '''指定された画像を安全に追加する(ファイルが存在しない場合はプレースホルダーを作成)'''
    if os.path.exists(image_path):
        try:
            return slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception as e:
            print(f"画像の追加に失敗しました: {e}")
            # 画像の追加に失敗した場合 代わりにテキストボックスを作成
            shape = slide.shapes.add_textbox(left, top, width, height)
            tf = shape.text_frame
            tf.text = f"[画像を表示できません: {os.path.basename(image_path)}]"
            return shape
    else:
        # 画像ファイルが存在しない場合 代わりにテキストボックスを作成
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.text = f"[画像ファイルが見つかりません: {os.path.basename(image_path)}]"
        return shape

# プレースホルダーに安全にアクセスするためのヘルパー関数
def get_placeholder_safe(slide, idx, default_title="", default_content=""):
    '''プレースホルダーに安全にアクセスし 存在しない場合はテキストボックスを作成'''
    try:
        return slide.placeholders[idx]
    except KeyError:
        # プレースホルダーが存在しない場合 テキストボックスを作成
        print(f"プレースホルダー {idx} が見つかりません. テキストボックスを作成します.")
        if idx == 0:  # タイトル用
            left, top, width, height = Inches(0.5), Inches(0.5), Inches(9), Inches(1.2)
            shape = slide.shapes.add_textbox(left, top, width, height)
            shape.text = default_title
            return shape
        elif idx == 1:  # 本文用
            left, top, width, height = Inches(0.5), Inches(2), Inches(9), Inches(4)
            shape = slide.shapes.add_textbox(left, top, width, height)
            shape.text = default_content
            return shape
        elif idx == 11:  # 特にプレースホルダー11のエラーに対応
            left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(3)
            shape = slide.shapes.add_textbox(left, top, width, height)
            shape.text = default_content
            return shape
        else:  # その他のインデックス
            left, top, width, height = Inches(1), Inches(1.5 + (idx * 0.5) % 5), Inches(8), Inches(1)
            shape = slide.shapes.add_textbox(left, top, width, height)
            shape.text = default_content
            return shape
"""

                            # 元のコードに必要なインポートが含まれているか確認し 不足していれば追加
                            modified_code = final_output
                            required_imports = [
                                "import os",
                                "from pptx.util import Inches, Pt",
                                "from pptx.enum.text import PP_ALIGN",
                                "from pptx.dml.color import RGBColor",
                                "from pptx.enum.shapes import MSO_SHAPE"
                            ]
                            
                            # 元のコードにインポート文がなければ追加
                            for import_line in required_imports:
                                if import_line not in modified_code:
                                    # コードの先頭に追加(既存のインポート文の後に)
                                    import_section_end = modified_code.find("\n\n", modified_code.find("import"))
                                    if import_section_end > 0:
                                        modified_code = modified_code[:import_section_end] + "\n" + import_line + modified_code[import_section_end:]
                                    else:
                                        # インポートセクションが見つからない場合は コードの先頭に追加
                                        modified_code = import_line + "\n\n" + modified_code
                            
                            # 先頭に安全関数を追加
                            modified_code = safe_placeholder_code + "\n\n" + modified_code
                            
                            # 画像ファイルのパスを安全な関数に置き換え
                            modified_code = modified_code.replace("slide.shapes.add_picture(", "add_image_safe(slide, ")
                            
                            # 直接アクセスを安全な関数に置き換え
                            modified_code = modified_code.replace("slide.placeholders[0]", "get_placeholder_safe(slide, 0)")
                            modified_code = modified_code.replace("slide.placeholders[1]", "get_placeholder_safe(slide, 1)")
                            # 特にプレースホルダー11に対応
                            modified_code = modified_code.replace("slide.placeholders[11]", "get_placeholder_safe(slide, 11)")
                            # 他のプレースホルダーインデックスにも対応（2から15まで）
                            for idx in range(2, 16):
                                modified_code = modified_code.replace(f"slide.placeholders[{idx}]", f"get_placeholder_safe(slide, {idx})")
                                
                            # フィル処理用のコードをさらに改善
                            modified_code = re.sub(
                                r'([a-zA-Z0-9_]+)\.fill\.fore_color\.rgb\s*=\s*([^;\n]+)',
                                r'set_fill_color_safe(\1.fill, \2)',
                                modified_code
                            )
                            
                            # フィル直接アクセスの場合も置き換え
                            modified_code = re.sub(
                                r'fill\.fore_color\.rgb\s*=\s*([^;\n]+)',
                                r'set_fill_color_safe(fill, \1)',
                                modified_code
                            )

                            # 修正したコードを保存
                            with open("workspace/output/create_pptx_safe.py", "w", encoding="utf-8") as f:
                                f.write(modified_code)
                            
                            # ユーザーに通知
                            st.info("プレースホルダーエラーを回避するために安全なバージョンも生成しました")
                        
                        # コードを実行
                        st.info("生成されたPythonコードを実行中...")
                        generation_successful = False # 成功フラグ
                        output_filename = None # 生成されたファイル名を格納

                        # execをラップするtryブロック開始
                        try:
                            # 安全なコードを実行する前に、必要なモジュールを明示的にインポート
                            import importlib
                            from pptx.util import Inches, Pt
                            from pptx.enum.text import PP_ALIGN
                            from pptx.dml.color import RGBColor
                            from pptx.enum.shapes import MSO_SHAPE
                            
                            # APIキーを環境変数に確実に設定
                            if api_provider == "OpenAI":
                                os.environ["OPENAI_API_KEY"] = api_key
                            else:  # Gemini
                                os.environ["GOOGLE_API_KEY"] = api_key
                            
                            # 実行環境に必要なモジュールを追加
                            exec_globals = {
                                "Inches": Inches,
                                "Pt": Pt,
                                "PP_ALIGN": PP_ALIGN,
                                "RGBColor": RGBColor,
                                "os": os,
                                "datetime": importlib.import_module('datetime'),
                                "random": importlib.import_module('random'),
                                "Presentation": importlib.import_module('pptx').Presentation,
                                "MSO_SHAPE": MSO_SHAPE,
                                # APIキーを直接渡す
                                "OPENAI_API_KEY": api_key if api_provider == "OpenAI" else "",
                                "GOOGLE_API_KEY": api_key if api_provider == "Google Gemini" else "",
                                # 安全な画像追加関数を追加
                                "add_image_safe": lambda slide, image_path, left=Inches(1), top=Inches(2), width=Inches(4), height=Inches(3): (
                                    slide.shapes.add_picture(image_path, left, top, width, height) 
                                    if os.path.exists(image_path) else 
                                    slide.shapes.add_textbox(left, top, width, height)
                                ),
                                # 塗りつぶし色設定用の安全関数を追加
                                "set_fill_color_safe": lambda fill, color: _set_fill_color(fill, color)
                            }
                            
                            # APIキーを使用したOpenAIインスタンス生成サポート
                            if api_provider == "OpenAI":
                                # OpenAIインポートを試行
                                try:
                                    from langchain_openai import ChatOpenAI, OpenAI
                                    exec_globals["ChatOpenAI"] = ChatOpenAI
                                    exec_globals["OpenAI"] = OpenAI
                                    
                                    # APIキーが設定されたOpenAIインスタンスを作成
                                    exec_globals["openai_llm"] = OpenAI(openai_api_key=api_key)
                                    exec_globals["openai_chat"] = ChatOpenAI(openai_api_key=api_key)
                                except ImportError:
                                    st.warning("langchain_openaiモジュールのインポートに失敗しました。修正されたコードで対応します。")
                            elif api_provider == "Google Gemini":
                                # Geminiインポートを試行
                                try:
                                    from langchain_google_genai import ChatGoogleGenerativeAI
                                    exec_globals["ChatGoogleGenerativeAI"] = ChatGoogleGenerativeAI
                                    
                                    # APIキーが設定されたGeminiインスタンスを作成
                                    exec_globals["gemini_chat"] = ChatGoogleGenerativeAI(model=model, google_api_key=api_key)
                                except ImportError:
                                    st.warning("langchain_google_genaiモジュールのインポートに失敗しました。修正されたコードで対応します。")
                                
                            # OpenAIのインポートと初期化コードパターンを探して置換
                            modified_code = re.sub(
                                r'from\s+langchain\s+import\s+OpenAI',
                                'from langchain_openai import OpenAI',
                                modified_code
                            )
                            
                            # OpenAIの初期化を修正
                            modified_code = re.sub(
                                r'OpenAI\(\s*\)',
                                f'OpenAI(openai_api_key="{api_key}")',
                                modified_code
                            )
                            modified_code = re.sub(
                                r'OpenAI\(\s*temperature\s*=\s*([^,\)]+)\s*\)',
                                f'OpenAI(temperature=\\1, openai_api_key="{api_key}")',
                                modified_code
                            )
                            
                            # fill 色設定用のヘルパー関数を追加
                            def _set_fill_color(fill, color):
                                if hasattr(fill, 'type') and fill.type == None:
                                    fill.solid()
                                if hasattr(fill, 'fore_color'):
                                    fill.fore_color.rgb = color
                                    return color
                                return None
                            
                            # _set_fill_color関数をexec_globalsに追加
                            exec_globals["_set_fill_color"] = _set_fill_color
                            
                            exec_locals = {}
                            
                            # 安全なコードを実行
                            exec(modified_code, exec_globals, exec_locals)
                            
                            generation_successful = True
                            st.success("プレゼンテーションの生成が完了しました！")

                            # 生成されたファイル名を取得
                            if 'output_filename' in exec_locals:
                                output_filename = exec_locals['output_filename']
                                if not os.path.isabs(output_filename):
                                    output_filename = os.path.abspath(os.path.join("workspace/output", output_filename))
                            if not output_filename or not os.path.exists(output_filename):
                                output_files = [f for f in os.listdir("workspace/output") if f.endswith(".pptx")]
                                if output_files:
                                    latest_file_path = max([os.path.join("workspace/output", f) for f in output_files], key=os.path.getmtime)
                                    output_filename = os.path.abspath(latest_file_path)

                        except KeyError as ke:
                            # プレースホルダーエラーの場合
                            if "no placeholder on this slide with idx" in str(ke):
                                st.warning("プレースホルダーに関するエラーが発生しました。")
                                
                                # 安全バージョンのコードがあれば、それを試す
                                safe_code_path = "workspace/output/create_pptx_safe.py"
                                if os.path.exists(safe_code_path):
                                    st.info("安全なバージョンのコードを試行します...")
                                    try:
                                        with open(safe_code_path, "r", encoding="utf-8") as f:
                                            safe_code = f.read()
                                        
                                        # 変数を確実に初期化
                                        if not 'safe_code' in locals() or safe_code is None or safe_code == "":
                                            # 安全なコードが読み込めない場合、元のコードを使用
                                            safe_code = final_output
                                        
                                        # 安全なコードを実行する前に、必要なモジュールを明示的にインポート
                                        import importlib
                                        from pptx.util import Inches, Pt
                                        from pptx.enum.text import PP_ALIGN
                                        from pptx.dml.color import RGBColor
                                        from pptx.enum.shapes import MSO_SHAPE
                                        
                                        # APIキーを環境変数に確実に設定
                                        if api_provider == "OpenAI":
                                            os.environ["OPENAI_API_KEY"] = api_key
                                        else:  # Gemini
                                            os.environ["GOOGLE_API_KEY"] = api_key
                                        
                                        # 実行環境に必要なモジュールを追加
                                        exec_globals = {
                                            "Inches": Inches,
                                            "Pt": Pt,
                                            "PP_ALIGN": PP_ALIGN,
                                            "RGBColor": RGBColor,
                                            "os": os,
                                            "datetime": importlib.import_module('datetime'),
                                            "random": importlib.import_module('random'),
                                            "Presentation": importlib.import_module('pptx').Presentation,
                                            "MSO_SHAPE": MSO_SHAPE,
                                            # APIキーを直接渡す
                                            "OPENAI_API_KEY": api_key if api_provider == "OpenAI" else "",
                                            "GOOGLE_API_KEY": api_key if api_provider == "Google Gemini" else "",
                                            # 安全な画像追加関数を追加
                                            "add_image_safe": lambda slide, image_path, left=Inches(1), top=Inches(2), width=Inches(4), height=Inches(3): (
                                                slide.shapes.add_picture(image_path, left, top, width, height) 
                                                if os.path.exists(image_path) else 
                                                slide.shapes.add_textbox(left, top, width, height)
                                            ),
                                            # 塗りつぶし色設定用の安全関数を追加
                                            "set_fill_color_safe": lambda fill, color: _set_fill_color(fill, color)
                                        }
                                        
                                        # APIキーを使用したOpenAIインスタンス生成サポート
                                        if api_provider == "OpenAI":
                                            # OpenAIインポートを試行
                                            try:
                                                from langchain_openai import ChatOpenAI, OpenAI
                                                exec_globals["ChatOpenAI"] = ChatOpenAI
                                                exec_globals["OpenAI"] = OpenAI
                                                
                                                # APIキーが設定されたOpenAIインスタンスを作成
                                                exec_globals["openai_llm"] = OpenAI(openai_api_key=api_key)
                                                exec_globals["openai_chat"] = ChatOpenAI(openai_api_key=api_key)
                                            except ImportError:
                                                st.warning("langchain_openaiモジュールのインポートに失敗しました。修正されたコードで対応します。")
                                        elif api_provider == "Google Gemini":
                                            # Geminiインポートを試行
                                            try:
                                                from langchain_google_genai import ChatGoogleGenerativeAI
                                                exec_globals["ChatGoogleGenerativeAI"] = ChatGoogleGenerativeAI
                                                
                                                # APIキーが設定されたGeminiインスタンスを作成
                                                exec_globals["gemini_chat"] = ChatGoogleGenerativeAI(model=model, google_api_key=api_key)
                                            except ImportError:
                                                st.warning("langchain_google_genaiモジュールのインポートに失敗しました。修正されたコードで対応します。")
                                        
                                        # OpenAIのインポートと初期化コードパターンを探して置換
                                        safe_code = re.sub(
                                            r'from\s+langchain\s+import\s+OpenAI',
                                            'from langchain_openai import OpenAI',
                                            safe_code
                                        )
                                        
                                        # OpenAIの初期化を修正
                                        safe_code = re.sub(
                                            r'OpenAI\(\s*\)',
                                            f'OpenAI(openai_api_key="{api_key}")',
                                            safe_code
                                        )
                                        safe_code = re.sub(
                                            r'OpenAI\(\s*temperature\s*=\s*([^,\)]+)\s*\)',
                                            f'OpenAI(temperature=\\1, openai_api_key="{api_key}")',
                                            safe_code
                                        )
                                        
                                        # fill 色設定用のヘルパー関数を追加
                                        def _set_fill_color(fill, color):
                                            if hasattr(fill, 'type') and fill.type == None:
                                                fill.solid()
                                            if hasattr(fill, 'fore_color'):
                                                fill.fore_color.rgb = color
                                                return color
                                            return None
                                        
                                        # _set_fill_color関数をexec_globalsに追加
                                        exec_globals["_set_fill_color"] = _set_fill_color
                                        
                                        exec_locals = {}
                                        
                                        # 安全なコードを実行
                                        exec(safe_code, exec_globals, exec_locals)
                                        
                                        generation_successful = True
                                        st.success("安全なコードでプレゼンテーションの生成が完了しました！")
                                        
                                        # 生成されたファイル名を取得
                                        if 'output_filename' in exec_locals:
                                            output_filename = exec_locals['output_filename']
                                            if not os.path.isabs(output_filename):
                                                output_filename = os.path.abspath(os.path.join("workspace/output", output_filename))
                                        if not output_filename or not os.path.exists(output_filename):
                                            output_files = [f for f in os.listdir("workspace/output") if f.endswith(".pptx")]
                                            if output_files:
                                                latest_file_path = max([os.path.join("workspace/output", f) for f in output_files], key=os.path.getmtime)
                                                output_filename = os.path.abspath(latest_file_path)
                                    
                                    except ImportError as import_err:
                                        st.error(f"必要なモジュールのインポートに失敗しました: {import_err}")
                                        st.info("組み込みの安全な生成関数を使用します...")
                                        try:
                                            output_filename = generate_safe_presentation()
                                            if output_filename and os.path.exists(output_filename):
                                                st.success("安全なプレゼンテーションの生成が完了しました！")
                                                st.info(f"生成されたプレゼンテーション: {output_filename}")
                                                generation_successful = True
                                        except Exception as safe_gen_err:
                                            st.error(f"安全なプレゼンテーション生成中にエラーが発生しました: {safe_gen_err}")
                                            import traceback
                                            st.error(traceback.format_exc())
                                    
                                    except Exception as safe_code_err:
                                        st.error(f"安全なコードの実行にも失敗しました: {safe_code_err}")
                                        # 組み込みの安全な生成関数を使用
                                        st.info("組み込みの安全な生成関数を使用します...")
                                        try:
                                            output_filename = generate_safe_presentation()
                                            
                                            if output_filename and os.path.exists(output_filename):
                                                st.success("安全なプレゼンテーションの生成が完了しました！")
                                                st.info(f"生成されたプレゼンテーション: {output_filename}")
                                                generation_successful = True
                                        except Exception as safe_gen_err:
                                            st.error(f"安全なプレゼンテーション生成中にエラーが発生しました: {safe_gen_err}")
                                            import traceback
                                            st.error(traceback.format_exc())
                                else:
                                    # 安全バージョンのコードがなければ、直接安全な生成関数を使用
                                    st.info("安全なプレゼンテーションを生成中...")
                                    try:
                                        output_filename = generate_safe_presentation()
                                        
                                        if output_filename and os.path.exists(output_filename):
                                            st.success("安全なプレゼンテーションの生成が完了しました！")
                                            st.info(f"生成されたプレゼンテーション: {output_filename}")
                                            generation_successful = True
                                    except Exception as safe_gen_err:
                                        st.error(f"安全なプレゼンテーション生成中にエラーが発生しました: {safe_gen_err}")
                                        import traceback
                                        st.error(traceback.format_exc())
                        except Exception as exec_e:
                            # exec実行中のその他のエラーの場合
                            st.error(f"コード実行中にエラーが発生しました: {str(exec_e)}")
                            import traceback
                            st.error(traceback.format_exc())

                        # ダウンロードボタンの表示 (成功した場合のみ) - tryブロック3の外側、ブロック2の内側
                        if generation_successful and output_filename:
                            # ファイルが存在するか最終確認
                            if os.path.exists(output_filename):
                                try:
                                    with open(output_filename, "rb") as f:
                                        st.download_button(
                                            label=f"{os.path.basename(output_filename)}をダウンロード",
                                            data=f,
                                            file_name=os.path.basename(output_filename),
                                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                            key=f"download-{os.path.basename(output_filename)}"
                                        )
                                except FileNotFoundError:
                                    st.error(f"ダウンロード用にファイルを開けませんでした: {output_filename}")
                                except Exception as download_err:
                                    st.error(f"ダウンロードボタンの表示中にエラーが発生しました: {download_err}")
                            else:
                                st.warning(f"生成されたはずのファイルが見つかりません: {output_filename}")
                        elif generation_successful:
                            st.warning("プレゼンテーションは生成されましたが、ダウンロード用のファイルが見つかりませんでした。")

                except requests.exceptions.RequestException as req_err: # APIコール関連のエラー
                    st.error(f"APIリクエスト中にエラーが発生しました: {req_err}")
                except Exception as e: # AI生成やファイル書き込み段階のエラー
                    st.error(f"プレゼンテーション生成プロセスでエラーが発生しました: {str(e)}")
                    error_msg = str(e).lower()
                    if any(err in error_msg for err in ["insufficient_quota", "quota exceeded", "rate_limit"]):
                        st.error(f"{api_provider} APIの制限に達しました。APIキーの制限を確認するか、別のAPIキーを使用してください。")
                    else:
                        st.error("詳細なエラー情報を確認するには、コンソールログを確認してください。")

    except Exception as main_e:
        st.error(f"メイン処理で予期せぬエラーが発生しました: {main_e}")
        import traceback
        st.error(traceback.format_exc())

    # 使用方法のガイド - ブロック1の外側で、main関数内
    with st.expander("使用方法"):
        st.write("""
        ### 使用方法
        1. テキストファイル (.txt, .docx, .md) をアップロードします. このファイルにはプレゼンテーションに含めたい内容を記述します.
        2. APIプロバイダー(OpenAIまたはGoogle Gemini)を選択します.
        3. 選択したプロバイダーのAPIキーを入力します.
        4. 使用するモデルを選択します.
        5. テンプレートファイル (.pptx) をアップロードします. これがプレゼンテーションのベースとなります.
        6. 必要に応じて 画像や動画ファイルをアップロードします.
        7. 「プレゼンテーションを生成」ボタンをクリックします.
        8. 生成されたプレゼンテーションをダウンロードします.
        
        ### 画像・動画の使用について
        テキストファイル内で画像や動画を使用したい場合は 以下のように指定してください:
        
        ```
        [画像: 説明テキスト]
        [動画: 説明テキスト]
        ```
        
        AIはこれらの指示に基づいて アップロードされた画像や動画ファイルを使用してプレゼンテーションを生成します.
        """)

if __name__ == "__main__":
    # ... (既存の依存関係チェックなど) ... 
    check_and_install_dependencies()
    ensure_directories()
    main() 